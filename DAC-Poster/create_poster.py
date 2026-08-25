#!/usr/bin/env python3
"""Build the CREDIT DAC Young Fellows poster from the supplied template."""

from __future__ import annotations

import csv
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.lines import Line2D
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
POSTER_DIR = Path(__file__).resolve().parent
ASSET_DIR = POSTER_DIR / "assets"
OUTPUT_DIR = POSTER_DIR / "output"
TEMPLATE = ROOT / "DACYF-PosterFormat.pptx"
RTX_RESULTS = ROOT / "evaluation/results/rtx5090/summary.csv"
H100_RESULTS = (
    ROOT
    / "modal_dsmem_evaluation/results/h100_20260713T184907Z/summary.csv"
)

WORKLOADS = [
    "layernorm_backward",
    "weighted_var_backward",
    "pearson_backward",
    "softmax_logits_backward",
    "lars_momentum",
    "rowwise_quant",
]
LABELS = {
    "layernorm_backward": "LayerNorm",
    "weighted_var_backward": "Weighted variance",
    "pearson_backward": "Pearson",
    "softmax_logits_backward": "Softmax logits",
    "lars_momentum": "LARS",
    "rowwise_quant": "Row quantization",
}

# Poster palette: DAC lime/orange, architecture blue, neutral graphite.
BLACK = "111111"
CHARCOAL = "33373B"
GRAY = "666D75"
LIGHT_GRAY = "E4E7EA"
PALE_GRAY = "F5F6F7"
WHITE = "FFFFFF"
DAC_LIME = "B7D52A"
DARK_LIME = "718C00"
PALE_LIME = "F4F8E7"
ORANGE = "F47A38"
PALE_ORANGE = "FFF2E9"
BLUE = "0072B2"
PALE_BLUE = "EAF4F9"
UW_RED = "C5050C"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def read_64k_speedups(path: Path) -> dict[str, float]:
    values: dict[str, float] = {}
    with path.open(newline="", encoding="utf-8") as source:
        for row in csv.DictReader(source):
            if int(row["cols"]) == 65536 and row["workload"] in WORKLOADS:
                values[row["workload"]] = float(row["dsmem_vs_best"])
    missing = set(WORKLOADS) - values.keys()
    if missing:
        raise RuntimeError(f"Missing 64K results in {path}: {sorted(missing)}")
    return values


def configure_plot_style() -> None:
    plt.rcParams.update(
        {
            "font.family": "DejaVu Sans",
            "font.size": 24,
            "axes.labelsize": 25,
            "xtick.labelsize": 21,
            "ytick.labelsize": 23,
            "axes.linewidth": 1.5,
            "pdf.fonttype": 42,
            "ps.fonttype": 42,
        }
    )


def plot_speedups() -> Path:
    rtx = read_64k_speedups(RTX_RESULTS)
    h100 = read_64k_speedups(H100_RESULTS)
    y = np.arange(len(WORKLOADS))
    figure, axis = plt.subplots(figsize=(10.2, 8.0), dpi=220)
    axis.axvspan(1.0, 2.55, color="#F4F8E7", zorder=0)
    axis.axvline(1.0, color="#222222", linewidth=2.0, linestyle=(0, (5, 4)))

    for index, workload in enumerate(WORKLOADS):
        axis.plot(
            [rtx[workload], h100[workload]],
            [index, index],
            color="#C8CDD2",
            linewidth=4,
            solid_capstyle="round",
            zorder=1,
        )
        axis.scatter(
            rtx[workload],
            index,
            s=205,
            color="#4D5156",
            edgecolor="white",
            linewidth=1.5,
            zorder=3,
        )
        axis.scatter(
            h100[workload],
            index,
            s=225,
            marker="s",
            color="#0072B2",
            edgecolor="white",
            linewidth=1.5,
            zorder=3,
        )
        axis.text(
            rtx[workload],
            index - 0.23,
            f"{rtx[workload]:.2f}",
            ha="center",
            va="center",
            fontsize=18,
            color="#34373A",
            fontweight="bold",
        )
        axis.text(
            h100[workload],
            index + 0.23,
            f"{h100[workload]:.2f}",
            ha="center",
            va="center",
            fontsize=18,
            color="#00639A",
            fontweight="bold",
        )

    axis.set_yticks(y, [LABELS[item] for item in WORKLOADS])
    axis.set_xlim(0.80, 2.55)
    axis.set_xticks([1.0, 1.5, 2.0, 2.5])
    axis.set_xlabel("Speedup over fastest baseline (×)", labelpad=10)
    axis.invert_yaxis()
    axis.spines[["top", "right", "left"]].set_visible(False)
    axis.tick_params(axis="y", length=0, pad=10)
    axis.tick_params(axis="x", length=7, width=1.5)
    axis.grid(False)
    legend = [
        Line2D(
            [0],
            [0],
            marker="o",
            color="none",
            markerfacecolor="#4D5156",
            markeredgecolor="white",
            markersize=13,
            label="RTX 5090",
        ),
        Line2D(
            [0],
            [0],
            marker="s",
            color="none",
            markerfacecolor="#0072B2",
            markeredgecolor="white",
            markersize=13,
            label="H100",
        ),
    ]
    axis.legend(
        handles=legend,
        loc="lower right",
        frameon=False,
        fontsize=22,
        ncol=2,
        columnspacing=1.0,
        handletextpad=0.3,
    )
    figure.tight_layout(pad=0.25)
    path = ASSET_DIR / "speedups_64k.png"
    figure.savefig(path, facecolor="white", bbox_inches="tight", pad_inches=0.06)
    plt.close(figure)
    return path


def plot_traffic() -> Path:
    reductions = {
        "LayerNorm": 50,
        "Weighted variance": 60,
        "Pearson": 33,
        "Softmax logits": 57,
        "LARS": 38,
        "Row quantization": 44,
    }
    labels = list(reductions)
    values = [reductions[label] for label in labels]
    y = np.arange(len(labels))
    figure, axis = plt.subplots(figsize=(10.0, 5.3), dpi=220)
    bars = axis.barh(y, values, height=0.57, color="#91AF14")
    for bar, value in zip(bars, values):
        axis.text(
            value + 1.4,
            bar.get_y() + bar.get_height() / 2,
            f"-{value}%",
            va="center",
            ha="left",
            fontsize=22,
            color="#4B5D00",
            fontweight="bold",
        )
    axis.set_yticks(y, labels)
    axis.set_xlim(0, 69)
    axis.invert_yaxis()
    axis.spines[:].set_visible(False)
    axis.tick_params(axis="y", length=0, pad=8)
    axis.set_xticks([])
    axis.grid(False)
    figure.tight_layout(pad=0.2)
    path = ASSET_DIR / "traffic_reduction.png"
    figure.savefig(path, facecolor="white", bbox_inches="tight", pad_inches=0.05)
    plt.close(figure)
    return path


def set_fill(shape, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_line(shape, color: str, width: float = 1.0, transparency: int = 0) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def remove_line(shape) -> None:
    shape.line.fill.background()


def add_text(
    slide,
    x: float,
    y: float,
    width: float,
    height: float,
    text: str,
    *,
    size: float = 26,
    color: str = BLACK,
    bold: bool = False,
    font: str = "Arial",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.02,
    fit: bool = False,
):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    lines = text.splitlines() or [""]
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    if fit:
        frame.fit_text(font_family=font, max_size=int(size))
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    fill: str = WHITE,
    line: str | None = None,
    line_width: float = 1.0,
    radius: bool = False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(width), Inches(height)
    )
    set_fill(shape, fill)
    if line is None:
        remove_line(shape)
    else:
        set_line(shape, line, line_width)
    return shape


def set_box_text(
    shape,
    text: str,
    *,
    size: float,
    color: str = BLACK,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
    valign=MSO_ANCHOR.MIDDLE,
    margin: float = 0.08,
) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    lines = text.splitlines() or [""]
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = line
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)


def add_section_header(slide, number: str, title: str, x: float, y: float, width: float):
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.72), Inches(0.72)
    )
    set_fill(circle, DAC_LIME)
    remove_line(circle)
    set_box_text(circle, number, size=27, color=BLACK, bold=True, margin=0)
    add_text(
        slide,
        x + 0.90,
        y + 0.03,
        width - 0.90,
        0.64,
        title.upper(),
        size=31,
        color=CHARCOAL,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x + 0.90),
        Inches(y + 0.68),
        Inches(width - 0.90),
        Inches(0.07),
    )
    set_fill(line, DAC_LIME)
    remove_line(line)


def add_down_arrow(slide, x: float, y: float, width: float, height: float, color: str):
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(width), Inches(height)
    )
    set_fill(arrow, color)
    remove_line(arrow)
    return arrow


def add_up_arrow(slide, x: float, y: float, width: float, height: float, color: str):
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.UP_ARROW, Inches(x), Inches(y), Inches(width), Inches(height)
    )
    set_fill(arrow, color)
    remove_line(arrow)
    return arrow


def add_right_arrow(slide, x: float, y: float, width: float, height: float, color: str):
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(width), Inches(height)
    )
    set_fill(arrow, color)
    remove_line(arrow)
    return arrow


def add_bullet(slide, x: float, y: float, width: float, text: str, color: str = ORANGE):
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y + 0.16), Inches(0.24), Inches(0.24)
    )
    set_fill(dot, color)
    remove_line(dot)
    add_text(slide, x + 0.42, y, width - 0.42, 0.66, text, size=25, color=CHARCOAL)


def replace_template_title(slide) -> None:
    title_shape = None
    footer_shape = None
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text.strip()
        if "Title of the poster" in text:
            title_shape = shape
        elif "logos of your university" in text:
            footer_shape = shape

    if title_shape is None:
        raise RuntimeError("Template title placeholder not found")

    frame = title_shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.word_wrap = True

    first = frame.paragraphs[0]
    first.alignment = PP_ALIGN.CENTER
    first.space_after = Pt(2)
    credit = first.add_run()
    credit.text = "CREDIT"
    credit.font.name = "Arial"
    credit.font.size = Pt(80)
    credit.font.bold = True
    credit.font.color.rgb = rgb(DARK_LIME)
    subtitle = first.add_run()
    subtitle.text = ": Cost-Guided DSMEM Tiling"
    subtitle.font.name = "Arial"
    subtitle.font.size = Pt(54)
    subtitle.font.bold = True
    subtitle.font.color.rgb = rgb(BLACK)

    second = frame.add_paragraph()
    second.alignment = PP_ALIGN.CENTER
    second.space_before = Pt(0)
    second.space_after = Pt(5)
    run = second.add_run()
    run.text = "for Wide Reduction-Reuse GPU Kernels"
    run.font.name = "Arial"
    run.font.size = Pt(52)
    run.font.bold = True
    run.font.color.rgb = rgb(BLACK)

    third = frame.add_paragraph()
    third.alignment = PP_ALIGN.CENTER
    third.space_before = Pt(0)
    third.space_after = Pt(0)
    run = third.add_run()
    run.text = (
        "Zhengxiong Li  |  Fellow ID: ______  |  Advisors: Tsung-Wei Huang & "
        "Umit Ogras  |  University of Wisconsin-Madison"
    )
    run.font.name = "Arial"
    run.font.size = Pt(25)
    run.font.color.rgb = rgb(CHARCOAL)

    if footer_shape is not None:
        footer_shape.text_frame.clear()


def draw_problem_panel(slide, x: float, y: float, width: float) -> None:
    add_section_header(slide, "1", "The problem", x, y, width)
    add_text(
        slide,
        x,
        y + 1.05,
        width,
        1.08,
        "Wide rows hit the memory wall",
        size=39,
        color=BLACK,
        bold=True,
    )
    add_text(
        slide,
        x,
        y + 2.02,
        width,
        0.75,
        "Conventional multi-pass schedule",
        size=24,
        color=GRAY,
        bold=True,
    )

    stage_y = y + 3.05
    stage_w = 2.72
    stage_gap = 0.50
    stages = [("REDUCE", "pass 1"), ("REDUCE", "pass 2"), ("OUTPUT", "pass 3")]
    for index, (label, sublabel) in enumerate(stages):
        sx = x + index * (stage_w + stage_gap)
        stage = add_rect(
            slide,
            sx,
            stage_y,
            stage_w,
            1.12,
            fill=PALE_GRAY,
            line=LIGHT_GRAY,
            line_width=1.8,
            radius=True,
        )
        set_box_text(stage, f"{label}\n{sublabel}", size=21, color=CHARCOAL, bold=True)

    hbm_y = stage_y + 2.33
    hbm = add_rect(
        slide,
        x,
        hbm_y,
        width - 0.15,
        1.12,
        fill=PALE_ORANGE,
        line=ORANGE,
        line_width=2.1,
    )
    set_box_text(hbm, "HBM: reload the same wide row", size=25, color=CHARCOAL, bold=True)
    for index in range(3):
        center = x + index * (stage_w + stage_gap) + stage_w / 2
        add_up_arrow(slide, center - 0.27, stage_y + 1.17, 0.54, 1.08, ORANGE)
        add_text(
            slide,
            center - 0.60,
            stage_y + 1.58,
            1.20,
            0.45,
            "read" if index == 0 else "reread",
            size=17,
            color=ORANGE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    bullet_y = hbm_y + 1.55
    add_bullet(slide, x, bullet_y, width, "One CTA cannot retain a wide row")
    add_bullet(slide, x, bullet_y + 0.75, width, "Later stages reread bulk data")
    add_bullet(slide, x, bullet_y + 1.50, width, "Naive clusters add barriers + remote traffic")

    cost_y = bullet_y + 2.63
    add_text(slide, x, cost_y, width, 0.62, "DSMEM IS NOT FREE", size=27, bold=True)
    latency = [
        ("Local SMEM", "30-34 cyc", 0.58, CHARCOAL),
        ("Remote DSMEM", "191-216 cyc", 1.75, BLUE),
        ("DRAM", "688-930 cyc", 4.20, ORANGE),
    ]
    for index, (label, value, bar_width, color) in enumerate(latency):
        row_y = cost_y + 0.78 + index * 0.83
        add_text(slide, x, row_y, 2.72, 0.52, label, size=21, color=GRAY)
        add_rect(slide, x + 2.78, row_y + 0.06, bar_width, 0.36, fill=color)
        add_text(
            slide,
            x + 7.18,
            row_y - 0.02,
            2.18,
            0.54,
            value,
            size=18,
            color=color,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )

    sync_y = cost_y + 3.43
    add_text(slide, x, sync_y, width, 0.55, "Cluster synchronization", size=22, color=GRAY)
    badge1 = add_rect(slide, x, sync_y + 0.63, 4.43, 0.92, fill=PALE_BLUE, line=BLUE, radius=True)
    set_box_text(badge1, "404 cycles  |  RTX 5090", size=20, color=BLUE, bold=True)
    badge2 = add_rect(slide, x + 4.70, sync_y + 0.63, 4.43, 0.92, fill=PALE_ORANGE, line=ORANGE, radius=True)
    set_box_text(badge2, "851 cycles  |  H100", size=20, color=ORANGE, bold=True)

    question_y = sync_y + 1.82
    question = add_rect(
        slide,
        x,
        question_y,
        width - 0.15,
        1.68,
        fill=PALE_ORANGE,
        line=ORANGE,
        line_width=2.3,
        radius=True,
    )
    set_box_text(
        question,
        "Can avoided rereads repay\ncluster overhead?",
        size=28,
        color=CHARCOAL,
        bold=True,
    )

    chart_y = question_y + 2.10
    add_text(slide, x, chart_y, width, 0.55, "WHY WIDTH HELPS", size=26, color=CHARCOAL, bold=True)
    axis_x = x + 0.80
    axis_y = chart_y + 3.45
    x_axis = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(axis_x),
        Inches(axis_y),
        Inches(x + width - 0.35),
        Inches(axis_y),
    )
    set_line(x_axis, CHARCOAL, 1.7)
    y_axis = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(axis_x),
        Inches(chart_y + 0.78),
        Inches(axis_x),
        Inches(axis_y),
    )
    set_line(y_axis, CHARCOAL, 1.7)
    saved = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(axis_x + 0.25),
        Inches(axis_y - 0.18),
        Inches(x + width - 0.45),
        Inches(chart_y + 0.92),
    )
    set_line(saved, DARK_LIME, 3.2)
    overhead_y = chart_y + 2.08
    overhead = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(axis_x + 0.25),
        Inches(overhead_y),
        Inches(x + width - 0.45),
        Inches(overhead_y),
    )
    set_line(overhead, ORANGE, 3.2)
    crossover_x = x + 5.18
    marker = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(crossover_x - 0.13),
        Inches(overhead_y - 0.13),
        Inches(0.26),
        Inches(0.26),
    )
    set_fill(marker, BLACK)
    remove_line(marker)
    crossover = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(crossover_x),
        Inches(overhead_y + 0.16),
        Inches(crossover_x),
        Inches(axis_y),
    )
    set_line(crossover, GRAY, 1.3)
    add_text(slide, x + 5.50, chart_y + 0.82, 3.55, 0.53, "avoided reread ∝ N", size=19, color=DARK_LIME, bold=True)
    add_text(slide, x + 5.50, overhead_y - 0.52, 3.55, 0.53, "cluster cost", size=19, color=ORANGE, bold=True)
    add_text(slide, crossover_x - 1.18, axis_y + 0.06, 2.36, 0.47, "crossover N*", size=18, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + width - 0.68, axis_y - 0.02, 0.55, 0.45, "N", size=20, color=CHARCOAL, bold=True, align=PP_ALIGN.RIGHT)


def draw_credit_panel(slide, x: float, y: float, width: float) -> None:
    highlight = add_rect(slide, x - 0.25, y - 0.28, width + 0.50, 22.38, fill=PALE_LIME)
    # Move the highlight behind later shapes but above the slide background.
    tree = slide.shapes._spTree
    tree.remove(highlight._element)
    tree.insert(2, highlight._element)

    add_section_header(slide, "2", "CREDIT", x, y, width)
    add_text(
        slide,
        x,
        y + 1.03,
        width,
        1.70,
        "KEEP BULK DATA LOCAL\nEXCHANGE ONLY SCALAR PARTIALS",
        size=36,
        color=DARK_LIME,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    row_y = y + 3.12
    segment_gap = 0.13
    segment_w = (width - 3 * segment_gap) / 4
    segment_colors = ["DDE9A4", "C8DCE8", "F5D8C5", "D9DDE1"]
    for index in range(4):
        sx = x + index * (segment_w + segment_gap)
        segment = add_rect(
            slide,
            sx,
            row_y,
            segment_w,
            0.95,
            fill=segment_colors[index],
            line=WHITE,
            line_width=1.0,
        )
        set_box_text(segment, f"ROW SLICE {index}", size=19, color=CHARCOAL, bold=True)
    add_text(
        slide,
        x,
        row_y - 0.43,
        width,
        0.35,
        "Partition one row across a thread-block cluster",
        size=20,
        color=GRAY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    local_y = row_y + 2.03
    scalar_y = local_y + 2.08
    for index in range(4):
        sx = x + index * (segment_w + segment_gap)
        add_down_arrow(slide, sx + segment_w / 2 - 0.23, row_y + 1.02, 0.46, 0.83, DARK_LIME)
        local = add_rect(
            slide,
            sx,
            local_y,
            segment_w,
            1.27,
            fill=WHITE,
            line=DARK_LIME,
            line_width=2.1,
            radius=True,
        )
        set_box_text(local, f"CTA {index}\nlocal SMEM", size=20, color=CHARCOAL, bold=True)
        add_down_arrow(slide, sx + segment_w / 2 - 0.22, local_y + 1.33, 0.44, 0.63, BLUE)
        token = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(sx + segment_w / 2 - 0.42),
            Inches(scalar_y),
            Inches(0.84),
            Inches(0.84),
        )
        set_fill(token, BLUE)
        remove_line(token)
        set_box_text(token, f"ρ{index}", size=18, color=WHITE, bold=True, margin=0)

    network_y = scalar_y + 1.17
    network = add_rect(
        slide,
        x + 0.40,
        network_y,
        width - 0.80,
        1.10,
        fill=PALE_BLUE,
        line=BLUE,
        line_width=2.2,
        radius=True,
    )
    set_box_text(
        network,
        "DSMEM replicated push: compact partials only",
        size=25,
        color=BLUE,
        bold=True,
    )
    for index in range(4):
        sx = x + index * (segment_w + segment_gap)
        add_down_arrow(slide, sx + segment_w / 2 - 0.20, scalar_y + 0.85, 0.40, 0.31, BLUE)

    output_y = network_y + 2.00
    for index in range(4):
        sx = x + index * (segment_w + segment_gap)
        add_down_arrow(slide, sx + segment_w / 2 - 0.23, network_y + 1.15, 0.46, 0.70, DARK_LIME)
        out = add_rect(
            slide,
            sx,
            output_y,
            segment_w,
            1.16,
            fill=WHITE,
            line=DARK_LIME,
            line_width=1.8,
            radius=True,
        )
        set_box_text(out, "reuse local\nslice → output", size=18, color=CHARCOAL, bold=True)

    complexity_y = output_y + 1.54
    callouts = [
        ("O(N/P)", "local bytes / CTA", DARK_LIME),
        ("O(P² Σq)", "remote bytes / row", BLUE),
        ("ZERO", "remote bulk reads", ORANGE),
    ]
    callout_gap = 0.22
    callout_w = (width - 2 * callout_gap) / 3
    for index, (headline, label, color) in enumerate(callouts):
        cx = x + index * (callout_w + callout_gap)
        box = add_rect(slide, cx, complexity_y, callout_w, 1.53, fill=WHITE, line=color, radius=True)
        frame = box.text_frame
        frame.clear()
        frame.margin_left = Inches(0.05)
        frame.margin_right = Inches(0.05)
        frame.margin_top = Inches(0.10)
        frame.margin_bottom = Inches(0.05)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = frame.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = headline
        r1.font.name = "Arial"
        r1.font.size = Pt(24)
        r1.font.bold = True
        r1.font.color.rgb = rgb(color)
        p2 = frame.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(0)
        r2 = p2.add_run()
        r2.text = label
        r2.font.name = "Arial"
        r2.font.size = Pt(17)
        r2.font.color.rgb = rgb(GRAY)

    model_y = complexity_y + 2.08
    add_text(
        slide,
        x,
        model_y,
        width,
        0.62,
        "COST-GUIDED GENERATION",
        size=27,
        color=CHARCOAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x,
        model_y + 0.69,
        width,
        0.72,
        "Predicted benefit =",
        size=24,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    terms = [
        ("saved\nreread", PALE_LIME, DARK_LIME, "+"),
        ("cluster\ncontrol", PALE_ORANGE, ORANGE, "-"),
        ("local\nreplay", PALE_GRAY, CHARCOAL, "-"),
        ("DSMEM\ntransport", PALE_BLUE, BLUE, "-"),
    ]
    term_gap = 0.17
    term_w = (width - 3 * term_gap) / 4
    term_y = model_y + 1.43
    for index, (label, fill, color, operator) in enumerate(terms):
        tx = x + index * (term_w + term_gap)
        term = add_rect(slide, tx, term_y, term_w, 1.25, fill=fill, line=color, radius=True)
        set_box_text(term, label, size=20, color=color, bold=True)
        if index > 0:
            operator_box = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(tx - 0.27),
                Inches(term_y + 0.42),
                Inches(0.38),
                Inches(0.38),
            )
            set_fill(operator_box, WHITE)
            remove_line(operator_box)
            set_box_text(operator_box, operator, size=23, color=CHARCOAL, bold=True, margin=0)

    decision_y = term_y + 1.55
    decision = add_rect(
        slide,
        x + 1.20,
        decision_y,
        width - 2.40,
        1.34,
        fill=DAC_LIME,
        line=DARK_LIME,
        line_width=2.0,
        radius=True,
    )
    set_box_text(decision, "GENERATE DSMEM KERNEL  iff  ΔT > 0", size=26, color=BLACK, bold=True)
    add_text(
        slide,
        x,
        decision_y + 1.57,
        width,
        0.95,
        "1 baseline timing + primitive profiles + static traffic counts\nNo DSMEM candidate timing",
        size=21,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )

    example_y = decision_y + 2.72
    example = add_rect(
        slide,
        x + 0.45,
        example_y,
        width - 0.90,
        1.95,
        fill=WHITE,
        line=DAC_LIME,
        line_width=2.3,
        radius=True,
    )
    frame = example.text_frame
    frame.clear()
    frame.margin_left = Inches(0.14)
    frame.margin_right = Inches(0.14)
    frame.margin_top = Inches(0.12)
    frame.margin_bottom = Inches(0.08)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = frame.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "LayerNorm at N = 65,536"
    r1.font.name = "Arial"
    r1.font.size = Pt(22)
    r1.font.bold = True
    r1.font.color.rgb = rgb(CHARCOAL)
    p2 = frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = "≈768 KiB HBM avoided  vs.  896 B DSMEM exchanged"
    r2.font.name = "Arial"
    r2.font.size = Pt(23)
    r2.font.bold = True
    r2.font.color.rgb = rgb(DARK_LIME)


def add_metric_badge(
    slide,
    x: float,
    y: float,
    width: float,
    headline: str,
    label: str,
    color: str,
    fill: str,
) -> None:
    badge = add_rect(slide, x, y, width, 1.62, fill=fill, line=color, line_width=1.8, radius=True)
    frame = badge.text_frame
    frame.clear()
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.06)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.04)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = frame.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = headline
    r1.font.name = "Arial"
    r1.font.size = Pt(31)
    r1.font.bold = True
    r1.font.color.rgb = rgb(color)
    p2 = frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(1)
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = "Arial"
    r2.font.size = Pt(17)
    r2.font.color.rgb = rgb(CHARCOAL)


def draw_evidence_panel(slide, x: float, y: float, width: float, speedup_plot: Path, traffic_plot: Path) -> None:
    add_section_header(slide, "3", "Cross-GPU evidence", x, y, width)
    add_text(
        slide,
        x,
        y + 1.02,
        width,
        0.83,
        "ALL 6 WORKLOADS WIN AT 64K",
        size=31,
        color=DARK_LIME,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x,
        y + 1.70,
        width,
        0.42,
        "best of torch.compile, Triton, and CUDA",
        size=18,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    slide.shapes.add_picture(
        str(speedup_plot), Inches(x - 0.05), Inches(y + 2.05), width=Inches(width + 0.10)
    )

    badge_y = y + 10.86
    add_metric_badge(slide, x, badge_y, 2.96, "1.466×", "RTX 5090 geomean", CHARCOAL, PALE_GRAY)
    add_metric_badge(slide, x + 3.18, badge_y, 2.96, "1.318×", "H100 geomean", BLUE, PALE_BLUE)
    add_metric_badge(slide, x + 6.36, badge_y, 2.96, "55 / 60", "CUDA profit signs", DARK_LIME, PALE_LIME)

    traffic_y = badge_y + 2.13
    add_text(
        slide,
        x,
        traffic_y,
        width,
        0.70,
        "MEASURED DRAM TRAFFIC ELIMINATED",
        size=25,
        color=CHARCOAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x,
        traffic_y + 0.55,
        width,
        0.52,
        "RTX 5090, N = 65,536",
        size=19,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    slide.shapes.add_picture(
        str(traffic_plot), Inches(x), Inches(traffic_y + 1.03), width=Inches(width)
    )

    model_y = traffic_y + 7.05
    model = add_rect(
        slide,
        x,
        model_y,
        width,
        2.25,
        fill=PALE_BLUE,
        line=BLUE,
        line_width=2.0,
        radius=True,
    )
    frame = model.text_frame
    frame.clear()
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.12)
    frame.margin_bottom = Inches(0.08)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = frame.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "MODEL VALIDATION"
    r1.font.name = "Arial"
    r1.font.size = Pt(23)
    r1.font.bold = True
    r1.font.color.rgb = rgb(CHARCOAL)
    p2 = frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = "90.0% RTX 5090  |  93.3% H100"
    r2.font.name = "Arial"
    r2.font.size = Pt(27)
    r2.font.bold = True
    r2.font.color.rgb = rgb(BLUE)
    p3 = frame.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(3)
    r3 = p3.add_run()
    r3.text = "Every crossover within one tested 2× width interval"
    r3.font.name = "Arial"
    r3.font.size = Pt(18)
    r3.font.color.rgb = rgb(CHARCOAL)


def draw_rule_band(slide, x: float, y: float, width: float) -> None:
    add_section_header(slide, "4", "Deployment rule", x, y, width)
    flow_y = y + 1.12
    flow_gap = 0.74
    flow_w = (width - 3 * flow_gap) / 4
    flow = [
        ("REREADS INPUT\nafter reduction?", PALE_GRAY, CHARCOAL),
        ("BULK STAYS\nowner-local?", PALE_LIME, DARK_LIME),
        ("REMOTE STATE\nstays compact?", PALE_BLUE, BLUE),
        ("ΔT > 0\nUSE DSMEM", DAC_LIME, BLACK),
    ]
    for index, (label, fill, color) in enumerate(flow):
        fx = x + index * (flow_w + flow_gap)
        box = add_rect(slide, fx, flow_y, flow_w, 1.42, fill=fill, line=color, line_width=1.8, radius=True)
        set_box_text(box, label, size=23, color=color, bold=True)
        if index < len(flow) - 1:
            add_right_arrow(slide, fx + flow_w + 0.10, flow_y + 0.49, 0.52, 0.44, GRAY)
    add_text(
        slide,
        x,
        flow_y + 1.52,
        width,
        0.45,
        "Any NO → keep the non-DSMEM baseline",
        size=19,
        color=ORANGE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    divider = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x + width / 2),
        Inches(y + 3.28),
        Inches(x + width / 2),
        Inches(y + 7.30),
    )
    set_line(divider, LIGHT_GRAY, 2.0)

    left_x = x
    right_x = x + width / 2 + 0.55
    half = width / 2 - 0.55
    add_text(slide, left_x, y + 3.20, half, 0.66, "USE DSMEM", size=28, color=DARK_LIME, bold=True)
    add_text(slide, right_x, y + 3.20, half, 0.66, "AVOID DSMEM", size=28, color=ORANGE, bold=True)

    use_items = [
        "Wide row + reduction + reuse",
        "Bulk slice stays owner-local",
        "Remote state stays compact",
    ]
    avoid_items = [
        "One-pass reduction",
        "Scan / selection",
        "Bulk halo exchange",
    ]
    for index, item in enumerate(use_items):
        add_bullet(slide, left_x, y + 3.96 + index * 0.70, half, item, color=DAC_LIME)
    for index, item in enumerate(avoid_items):
        add_bullet(slide, right_x, y + 3.96 + index * 0.70, half, item, color=ORANGE)

    add_text(
        slide,
        left_x,
        y + 6.25,
        half,
        0.58,
        "LayerNorm  |  softmax backward  |  pairwise statistics  |  quantization",
        size=19,
        color=GRAY,
        bold=True,
    )
    add_text(
        slide,
        right_x,
        y + 6.25,
        half,
        0.58,
        "Width amortizes overhead; width alone is not predictive.",
        size=19,
        color=GRAY,
        bold=True,
    )

    add_text(
        slide,
        x,
        y + 7.58,
        width,
        0.62,
        "Wide rows amortize overhead; width alone is not a selection rule.",
        size=26,
        color=GRAY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    takeaway = add_rect(slide, x, y + 9.18, width, 1.36, fill=CHARCOAL)
    set_box_text(
        takeaway,
        "DSMEM is a selective reduction-reuse mechanism, not a general fusion fabric.",
        size=31,
        color=WHITE,
        bold=True,
    )


def add_footer(slide) -> None:
    add_text(
        slide,
        1.08,
        45.78,
        11.8,
        1.22,
        "UNIVERSITY OF\nWISCONSIN-MADISON",
        size=23,
        color=UW_RED,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        12.0,
        45.95,
        12.0,
        0.75,
        "zhengxiong.li@wisc.edu",
        size=23,
        color=CHARCOAL,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        24.3,
        45.82,
        10.6,
        0.98,
        "RTX 5090 + NVIDIA H100\nCUDA 13  |  FP32",
        size=20,
        color=GRAY,
        bold=True,
        align=PP_ALIGN.RIGHT,
        valign=MSO_ANCHOR.MIDDLE,
    )


def build_poster() -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    configure_plot_style()
    speedup_plot = plot_speedups()
    traffic_plot = plot_traffic()

    presentation = Presentation(str(TEMPLATE))
    if len(presentation.slides) != 1:
        raise RuntimeError("Expected a one-slide DAC poster template")
    if round(presentation.slide_width / 914400, 2) != 36.0 or round(
        presentation.slide_height / 914400, 2
    ) != 48.0:
        raise RuntimeError("The template is not the required 36 x 48 inch format")

    slide = presentation.slides[0]
    replace_template_title(slide)

    left_x, left_w = 1.03, 9.55
    center_x, center_w = 11.18, 13.62
    right_x, right_w = 25.38, 9.56
    body_y = 9.64
    draw_problem_panel(slide, left_x, body_y, left_w)
    draw_credit_panel(slide, center_x, body_y, center_w)
    draw_evidence_panel(slide, right_x, body_y, right_w, speedup_plot, traffic_plot)
    draw_rule_band(slide, 1.03, 32.42, 33.91)
    add_footer(slide)

    output = OUTPUT_DIR / "CREDIT_DAC_Young_Fellows_Poster.pptx"
    presentation.save(output)
    return output


if __name__ == "__main__":
    output_path = build_poster()
    print(output_path)
